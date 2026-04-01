"""PowerPoint file parser (.pptx)."""

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches


def parse_pptx(file_path: Path) -> list[dict[str, Any]]:
    """
    Parse PowerPoint presentations, extracting text and speaker notes.

    Args:
        file_path: Path to the PowerPoint file.

    Returns:
        List of chunks with text and metadata per slide.
    """
    try:
        prs = Presentation(file_path)
    except Exception as e:
        raise ValueError(f"Failed to open PowerPoint: {e}") from e

    chunks = []

    for slide_num, slide in enumerate(prs.slides, 1):
        title = _get_slide_title(slide)
        content = _get_slide_content(slide)
        notes = _get_slide_notes(slide)

        # Build slide text
        text_parts = [f"Slide {slide_num}: {title}"]

        if content:
            text_parts.append(content)

        if notes:
            text_parts.append(f"Notes: {notes}")

        text = "\n\n".join(text_parts)

        if text.strip():
            chunks.append(
                {
                    "text": text,
                    "metadata": {
                        "source": file_path.name,
                        "file_type": "pptx",
                        "page_or_sheet_or_slide": slide_num,
                        "slide_title": title,
                    },
                }
            )

    return chunks


def _get_slide_title(slide) -> str:
    """Extract slide title."""
    if slide.shapes.title:
        return slide.shapes.title.text.strip()

    # Try to find title placeholder
    for shape in slide.shapes:
        if shape.is_placeholder:
            if hasattr(shape, "placeholder_format"):
                if shape.placeholder_format.type == 1:  # Title placeholder
                    if hasattr(shape, "text"):
                        return shape.text.strip()

    return "Untitled Slide"


def _get_slide_content(slide) -> str:
    """Extract all text content from slide shapes."""
    content_parts = []

    for shape in slide.shapes:
        # Skip title shape
        if shape == slide.shapes.title:
            continue

        if hasattr(shape, "text") and shape.text.strip():
            content_parts.append(shape.text.strip())

        # Handle tables
        if shape.has_table:
            table_text = _extract_table_text(shape.table)
            if table_text:
                content_parts.append(table_text)

    return "\n\n".join(content_parts)


def _get_slide_notes(slide) -> str:
    """Extract speaker notes from slide."""
    if slide.has_notes_slide:
        notes_slide = slide.notes_slide
        if notes_slide.notes_text_frame:
            return notes_slide.notes_text_frame.text.strip()
    return ""


def _extract_table_text(table) -> str:
    """Extract text from a PowerPoint table."""
    rows = []

    for row in table.rows:
        cells = [cell.text.strip() for cell in row.cells]
        rows.append(" | ".join(cells))

    if not rows:
        return ""

    return "Table:\n" + "\n".join(rows)
